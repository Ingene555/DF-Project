""" 
This library will be used for conversion into
A class per file type will be defined
"""

from PIL import Image
from toIMG import VideoToImg, ModifyToImg
import mainFunctions as mf
from datetime import datetime as dt
import os
import pdfplumber as pdfp
from docx import Document as doc
from pptx import Presentation as pre
import win32com.client as client32
import pandas as pd
import time

TEMPCACHES = "tempcaches"
MULTIPLE = "multiple"
ONE = "one"


class TxtToTxt:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.stop = None
        self.progressShow = None

    def toTXT(self, path, method=MULTIPLE):
        path = path if type(path) == list else [path]
        _ = list()
        if method == MULTIPLE:
            for xx, ww in enumerate(self.file):
                if not self.stop:
                    try:
                        with open(ww, 'r', encoding="utf-8") as f1:
                            with open(path[xx], 'w', encoding="utf-8") as f2:
                                f2.write(f1.read())
                                _.append(mf.setReturn(True, [type(path[xx])], [path[xx]]))
                    except Exception as e:
                        _.append(mf.setReturn(False, [type(e)], [e]))
                    if self.progressShow:
                        self.progressShow(int(xx * 100 / len(self.file)))
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        elif method == ONE:
            a = ""
            for ww in self.file:
                if not self.stop:
                    try:
                        with open(ww, 'r', encoding="utf-8") as f:
                            a += f.read()
                    except Exception as e:
                        print(e)
                    a += "\n"*10
                else:
                    break
            try:
                for ww in path:
                    if not self.stop:
                        try:
                            with open(ww, 'w', encoding="utf-8") as f:
                                f.write(a)
                            _.append(mf.setReturn(True, [type(ww)], [ww]))
                        except Exception as e:
                            _.append(mf.setReturn(False, [type(e)], [e]))
                        if self.progressShow:
                            self.progressShow(int(path.index(ww) * 100 / len(path)))
                    else:
                        for i in _:
                            try:
                                os.remove(mf.Return(i).getValues(0))
                            except:pass
                        _=[]
                        break
            except Exception as e:
                for ww in path:
                    _.append(mf.setReturn(False, [type(e)], [e]))
                    if self.progressShow:
                        self.progressShow(int(path.index(ww) * 100 / len(path)))
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        

class ImgToTxt:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None

    def Convert(self, pathF, pathT, char="@%#*+=-:. ", modify={}, convert=False):
        pathT
        _ = list()
        for ww in range(len(pathF)):
            type(ww)
            _.append(os.path.join(TEMPCACHES, f"tc-imgxt-{ww}-{dt.now()}.txt".replace(":", "")))
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-imgxt-{dt.now()}.txt".replace(":", "")))
        __ = list()
        self.toKill = ModifyToImg(pathF)
        pathF = self.toKill.toIMG([f"{TEMPCACHES}/tc-{ww}-{dt.now()}.png".replace(":", "") for ww in range(len(pathF))], modify) if convert else [mf.setReturn(True, [type(ww)], [ww]) for ww in pathF]
        fx = 0

        def pixel_to_ascii(pixel_value):
            return ascii_chars[pixel_value // 32]

        for xx, ww in enumerate(pathF):
            if not self.stop:
                fy = 0
                if mf.Return(ww).getSucced():
                    try:
                        with Image.open(mf.Return(ww).getValues(0)) as img:
                            img = img.convert('L')
                            width, height = img.size
                            pixels = img.load()
                            ascii_chars = char
                            with open(_[xx], 'w', encoding="utf-8") as f:
                                for y in range(height):
                                    for x in range(width):
                                        gray_value = pixels[x, y]
                                        ascii_char = pixel_to_ascii(gray_value)
                                        f.write(ascii_char)
                                    f.write('\n')
                                    fy = int(y * 100 / height)
                                    if self.progressShow:
                                        self.progressShow(fy)
                            __.append(mf.setReturn(True, [type(_[xx])], [_[xx]]))
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:
                            print(e)
                    except Exception as e:
                        __.append(mf.setReturn(False, [type(e)], [e]))
                else:
                    __.append(ww)
                fy = 100
                fx = int(xx * 100 / len(pathF))
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __

    def toTXT(self, path, char="@%#*+=-:. ", method=MULTIPLE, modify={}, convert=True):
        temp = self.Convert(self.file, [], char, modify, convert)
        path = path if type(path) == list else [path]
        _ = list()
        if method == MULTIPLE:
            for xx, ww in enumerate(temp):
                if not self.stop:
                    if mf.Return(ww).getSucced():
                        with open(mf.Return(ww).getValues(0), 'r', encoding="utf-8") as f1:
                            with open(path[xx], 'w', encoding="utf-8") as f2:
                                f2.write(f1.read())
                                _.append(mf.setReturn(True, [type(path[xx])], [path[xx]]))
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:
                            print(e)
                    else:
                        _.append(ww)
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        elif method == ONE:
            a = ""
            for xx, ww in enumerate(temp):
                if not self.stop:
                    if mf.Return(ww).getSucced():
                        with open(mf.Return(ww).getValues(0), 'r', encoding="utf-8") as f:
                            a += f.read()
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:
                            print(e)
                    else:
                        print(ww)
                    a += "\n"*10
                else:
                    break
            for ww in path:
                if not self.stop:
                    try:
                        with open(ww, 'w', encoding="utf-8") as f:
                            f.write(a)
                            _.append(mf.setReturn(True, [type(ww)], [ww]))
                    except Exception as e:
                        _.append(mf.setReturn(False, [type(e)], [e]))
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        else:
            _.append(mf.setReturn(False, [type(None)], [None]))  
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)
        

class VideoToTxt:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None

    def Convert(self, pathF, pathT, interval=None, char="@%#*+=-:. ", modify={}, method=MULTIPLE):
        pathT
        _ = list()
        for ww in range(len(pathF)):
            type(ww)
            _.append(os.path.join(TEMPCACHES, f"tc-videoxt-{ww}-{dt.now()}.txt".replace(":", "")).replace('\\', '/'))
            os.mkdir(_[-1])
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-{dt.now()}.txt".replace(":", "")).replace('\\', '/'))
            os.mkdir(_[-1])
        __ = list()
        self.toKill = VideoToImg(pathF)
        modImg = self.toKill.toIMG([f"{TEMPCACHES}/tc-{ww}-{dt.now()}".replace(":", "") for ww in range(len(pathF))], ["png" for ww in range(len(pathF))], interval, modify)
        for xx, ww in enumerate(modImg):
            if not self.stop:
                self.toKill = ImgToTxt([mf.Return(ww[yy]).getValues(0).replace("\\", "/") for yy in range(len(ww))])
                __.append(self.toKill.toTXT([f"{_[xx]}/Frame{yy}.txt" for yy in range(len(ww))], char, method, {}, False))
                try:[os.remove(mf.Return(ww[yy]).getValues(0)) for yy in range(len(ww))]
                except Exception as e:print(e)
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        return __
    
    def toTXT(self, path, interval=None, char="@%#*+=-:. ", modify={}, method=MULTIPLE):
        temp = self.Convert(self.file, [], interval, char, modify, method)
        path = path if type(path) == list else [path]
        _ = list()
        for xx, ww in enumerate(temp):
            if not self.stop:
                __ = list()
                if not os.path.exists(path[xx]):os.mkdir(path[xx])
                if method == MULTIPLE:
                    for yy, zz in enumerate(ww):
                        if not self.stop:
                            if mf.Return(zz).getSucced():
                                with open(mf.Return(zz).getValues(0), 'r', encoding="utf-8") as f1:
                                    with open(f"{path[xx]}/Frame{yy}.txt".replace("\\", "/"), 'w', encoding="utf-8") as f2:
                                        f2.write(f1.read())
                                        __.append(mf.setReturn(True, [type(f"{path[xx]}/Frame{yy}.txt".replace("\\", "/"))], [f"{path[xx]}/Frame{yy}.txt".replace("\\", "/")]))
                                try:os.remove(mf.Return(zz).getValues(0))
                                except Exception as e:
                                    print(e)
                            else:
                                __.append(zz)
                        else:
                            break
                elif method == ONE:
                    a = ""
                    for yy, zz in enumerate(ww):
                        if not self.stop:
                            if mf.Return(zz).getSucced():
                                with open(mf.Return(zz).getValues(0), 'r', encoding="utf-8") as f:
                                    a += f.read()
                                try:os.remove(mf.Return(zz).getValues(0))
                                except Exception as e:
                                    print(e)
                            else:
                                print(zz)
                            a += "\n"*10
                        else:
                            break
                    with open(f"{path[xx]}/Frame.txt".replace("\\", "/"), 'w', encoding="utf-8") as f:
                        f.write(a)
                        __.append(mf.setReturn(True, [type(f"{path[xx]}/Frame.txt".replace("\\", "/"))], [f"{path[xx]}/Frame.txt".replace("\\", "/")]))
                else:
                    __.append(mf.setReturn(False, [type(None)], [None]))
                _.append(__)
            else:
                for i in _:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                _=[]
                break
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)
        

class PdfToTxt:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None

    def Convert(self, pathF, pathT):
        _ = list()
        for ww in range(len(pathF)):
            _.append(os.path.join(TEMPCACHES, f"tc-pdfxt-{ww}-{dt.now()}.txt".replace(":", "")))
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-pdfxt-{dt.now()}.txt".replace(":", "")))
        __ = list()
        x = 0
        for ww in range(len(pathF)):
            if not self.stop:
                y = 0
                try:
                    with pdfp.open(pathF[ww]) as pdf:
                        text = ""
                        for xx, page in enumerate(pdf.pages):
                            text += page.extract_text() + "\n"
                            y = int(xx * 100 / len(pdf.pages))
                            if self.progressShow:
                                self.progressShow(y)
                        with open(_[ww], 'w', encoding="utf-8") as f:
                            f.write(text)
                            __.append(mf.setReturn(True, [type(_[ww])], [_[ww]]))
                except Exception as e:
                    __.append(mf.setReturn(True, [type(e)], [e]))
                x = int(ww * 100 / len(pathF))
                y = 100
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __

    def toTXT(self, path, method=MULTIPLE):
        temp = self.Convert(self.file, [])
        path = path if type(path) == list else [path]
        _ = list()
        if method == MULTIPLE:
            for xx, ww in enumerate(temp):
                if not self.stop:
                    if mf.Return(ww).getSucced():
                        with open(mf.Return(ww).getValues(0), 'r', encoding="utf-8") as f1:
                            with open(path[xx], 'w', encoding="utf-8") as f2:
                                f2.write(f1.read())
                                _.append(mf.setReturn(True, [type(path[xx])], [path[xx]]))
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:print(e)
                    else:
                        _.append(ww)
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        elif method == ONE:
            a = ""
            for ww in temp:
                if not self.stop:
                    if mf.Return(ww).getSucced():
                        with open(mf.Return(ww).getValues(0), 'r', encoding="utf-8") as f:
                            a += f.read()
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:print(e)
                    else:print(ww)
                    a += "\n"*10
                else:
                    break
            for ww in range(len(path)):
                if not self.stop:
                    try:
                        with open(path[ww], 'w', encoding="utf-8") as f:
                            f.write(a)
                            _.append(mf.setReturn(True, [type(path[ww])], [path[ww]]))
                    except Exception as e:
                        _.append(mf.setReturn(False, [type(e)], [e]))
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        else:
            _.append(mf.setReturn(False, [type(None)], [None])) 
        if self.progressShow:
            self.progressShow(100, True, self.stop, _) 
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)
        

class WordToTxt:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None

    def Convert(self, pathF, pathT):
        pathT
        _ = list()
        for ww in range(len(pathF)):
            type(ww)
            _.append(os.path.join(TEMPCACHES, f"tc-wordxt-{ww}-{dt.now()}.txt".replace(":", "")))
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-wprdxt-{dt.now()}.txt".replace(":", "")))
        __ = list()
        for xx, ww in enumerate(pathF):
            if not self.stop:
                try:
                    if mf.Return(mf.getExtension(ww)).getValues(0) in ["docx", "dotx"]:
                        text = self.x(ww)
                    else:text = self.o(ww)
                    with open(_[xx], 'w', encoding="utf-8") as f:
                        f.write(text)
                        __.append(mf.setReturn(True, [type(_[xx])], [_[xx]]))
                except Exception as e:
                    __.append(mf.setReturn(False, [type(e)], [e]))
                if self.progressShow:
                    self.progressShow(int(xx * 100 / len(pathF)))
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __ 

    def toTXT(self, path, method=MULTIPLE):
        temp = self.Convert(self.file, [])
        path = path if type(path) == list else [path]
        _ = list()
        if method == MULTIPLE:
            for xx, ww in enumerate(temp):
                if not self.stop:
                    if mf.Return(ww).getSucced():
                        with open(mf.Return(ww).getValues(0), 'r', encoding="utf-8") as f1:
                            with open(path[xx], 'w', encoding="utf-8") as f2:
                                f2.write(f1.read())
                                _.append(mf.setReturn(True, [type(path[xx])], [path[xx]]))
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:print(e)
                    else:
                        _.append(ww)
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        elif method == ONE:
            a = ""
            for ww in temp:
                if not self.stop:
                    if mf.Return(ww).getSucced():
                        with open(mf.Return(ww).getValues(0), 'r', encoding="utf-8") as f:
                            a += f.read()
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:print(e)
                    else:print(ww)
                    a += "\n"*10
                else: 
                    break
            for ww in range(len(path)):
                if not self.stop:
                    try:
                        with open(path[ww], 'w', encoding="utf-8") as f:
                            f.write(a)
                            _.append(mf.setReturn(True, [type(path[ww])], [path[ww]]))
                    except Exception as e:
                        _.append(mf.setReturn(False, [type(e)], [e]))
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        else:
            _.append(mf.setReturn(False, [type(None)], [None])) 
        if self.progressShow:
            self.progressShow(100, True, self.stop, _) 
        return _
    
    def x(self, docx):
        text = ""
        for para in doc(docx).paragraphs:
            text += para.text + "\n"
        return text

    def o(self, doc):
        doc = os.path.normpath(os.path.abspath(doc))
        word = client32.Dispatch("Word.Application")
        docy = word.Documents.Open(doc)
        doco = docy.Content.Text
        docy.Close()
        word.Quit()
        time.sleep(.1)
        return doco

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)
        

class ExcelToTxt:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None

    def Convert(self, pathF, pathT):
        pathT
        _ = list()
        for ww in range(len(pathF)):
            type(ww)
            _.append(os.path.join(TEMPCACHES, f"tc-excelxt-{ww}-{dt.now()}.txt".replace(":", "")))
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-excelxt-{dt.now()}.txt".replace(":", "")))
        __ = list()
        for xx, ww in enumerate(pathF):
            if not self.stop:
                try:
                    dm = None
                    if ww.endswith(".xlsb"):
                        df = pd.read_excel(ww, engine="pyxlsb")
                        dm = True
                    elif ww.endswith((".xls", ".xlsx", ".xlsm", ".xltx", ".xlt")):
                        df = pd.read_excel(ww)
                        dm = True
                    elif ww.endswith(".csv"):
                        try:
                            df = pd.read_csv(ww, encoding="utf-8")
                            dm = True
                        except:
                            df = pd.read_csv(ww, encoding="latin1")
                            dm = True
                    if dm:
                        df.to_csv(_[xx], sep="\t", index=False, encoding="utf-8")
                    else:
                        df = pd.read_excel(ww, sheet_name=None)
                        txt_content = ""
                        for sheet_name, data in df.items():
                            txt_content += f"Feuille : {sheet_name}\n"
                            txt_content += data.to_csv(sep="\t", index=False, header=True)
                            txt_content += "\n\n"
                        with open(_[xx], "w", encoding="utf-8") as txt_file:
                            txt_file.write(txt_content)
                    __.append(mf.setReturn(True, [type(_[xx])], [_[xx]])) 
                except Exception as e:
                    __.append(mf.setReturn(False, [type(e)], [e]))
                if self.progressShow:
                    self.progressShow(int(xx * 100 / len(pathF)))
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __

    def toTXT(self, path, method=MULTIPLE):
        temp = self.Convert(self.file, [])
        path = path if type(path) == list else [path]
        _ = list()
        if method == MULTIPLE:
            for xx, ww in enumerate(temp):
                if not self.stop:
                    if mf.Return(ww).getSucced():
                        with open(mf.Return(ww).getValues(0), 'r', encoding="utf-8") as f1:
                            with open(path[xx], 'w', encoding="utf-8") as f2:
                                f2.write(f1.read())
                                _.append(mf.setReturn(True, [type(path[xx])], [path[xx]]))
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:print(e)
                    else:
                        _.append(ww)
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        elif method == ONE:
            a = ""
            for ww in temp:
                if not self.stop:
                    if mf.Return(ww).getSucced():
                        with open(mf.Return(ww).getValues(0), 'r', encoding="utf-8") as f:
                            a += f.read()
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:print(e)
                    else:print(ww)
                    a += "\n"*10
                else: 
                    break
            for ww in range(len(path)):
                if not self.stop:
                    try:
                        with open(path[ww], 'w', encoding="utf-8") as f:
                            f.write(a)
                            _.append(mf.setReturn(True, [type(path[ww])], [path[ww]]))
                    except Exception as e:
                        _.append(mf.setReturn(False, [type(e)], [e]))
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        else:
            _.append(mf.setReturn(False, [type(None)], [None]))  
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)
        

class PpToTxt:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None

    def Convert(self, pathF, pathT):
        pathT
        _ = list()
        for ww in range(len(pathF)):
            type(ww)
            _.append(os.path.join(TEMPCACHES, f"tc-ppxt-{ww}-{dt.now()}.txt".replace(":", "")))
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-ppxt-{dt.now()}.txt".replace(":", "")))
        __ = list()
        for xx, ww in enumerate(pathF):
            if not self.stop:
                try:
                    file_extension = os.path.splitext(ww)[1].lower()
                    removable = False
                    if file_extension in ['.ppt', '.pps', '.ppsx']:
                        pptx_file = os.path.normpath(os.path.abspath(_[xx].replace(".txt", '.pptx')))
                        powerpoint = client32.Dispatch("PowerPoint.Application")
                        presentation = powerpoint.Presentations.Open(os.path.normpath(os.path.abspath(ww)))
                        presentation.SaveAs(pptx_file, 24)
                        presentation.Close()
                        powerpoint.Quit()
                        time.sleep(.1)
                        file_path = pptx_file
                        removable = True
                    else:
                        file_path = os.path.normpath(os.path.abspath(ww))
                    presentation = pre(file_path)
                    extracted_text = []
                    for slidex, slide in enumerate(presentation.slides):
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                extracted_text.append(shape.text)
                        if self.progressShow:
                            self.progressShow(int(slidex * 100 / len(presentation.slides)))
                    text = "\n".join(extracted_text)
                    with open(_[xx], 'w', encoding="utf-8") as f:f.write(text)
                    if removable:
                        try:os.remove(file_path)
                        except Exception as e:print(e)
                    __.append(mf.setReturn(True, [type(_[xx])], [_[xx]]))
                except Exception as e:
                    __.append(mf.setReturn(False, [type(e)], [e]))
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __

    def toTXT(self, path, method=MULTIPLE):
        temp = self.Convert(self.file, [])
        path = path if type(path) == list else [path]
        _ = list()
        if method == MULTIPLE:
            for xx, ww in enumerate(temp):
                if not self.stop:
                    if mf.Return(ww).getSucced():
                        with open(mf.Return(ww).getValues(0), 'r', encoding="utf-8") as f1:
                            with open(path[xx], 'w', encoding="utf-8") as f2:
                                f2.write(f1.read())
                                _.append(mf.setReturn(True, [type(path[xx])], [path[xx]]))
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:print(e)
                    else:
                        _.append(ww)
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        elif method == ONE:
            a = ""
            for ww in temp:
                if not self.stop:
                    if mf.Return(ww).getSucced():
                        with open(mf.Return(ww).getValues(0), 'r', encoding="utf-8") as f:
                            a += f.read()
                        try:os.remove(mf.Return(ww).getValues(0))
                        except Exception as e:print(e)
                    else:print(ww)
                    a += "\n"*10
                else:
                    break
            for ww in range(len(path)):
                if not self.stop:
                    try:
                        with open(path[ww], 'w', encoding="utf-8") as f:
                            f.write(a)
                            _.append(mf.setReturn(True, [type(path[ww])], [path[ww]]))
                    except Exception as e:
                        _.append(mf.setReturn(False, [type(e)], [e]))
                else:
                    for i in _:
                        try:
                            os.remove(mf.Return(i).getValues(0))
                        except:pass
                    _=[]
                    break
        else:
            _.append(mf.setReturn(False, [type(None)], [None]))  
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)



