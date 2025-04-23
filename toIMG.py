""" 
This library will be used for conversion into image
A class per file type will be defined
"""

from PIL import Image, ImageDraw, ImageFont, ImageEnhance
import mainFunctions as mf
from datetime import datetime as dt
import os
from wand.image import Image as WandImage
import win32api
import cv2
import fitz
from docx2pdf import convert as cwp
from pptx import Presentation 
import qrcode
import numpy as np
import win32com.client as client32
import time

TEMPCACHES = 'tempcaches'
ALL = "all"


def wtw(input_file, output_file):
    word = client32.Dispatch("Word.Application")
    word.Visible = False
    doc = word.Documents.Open(os.path.abspath(input_file))
    doc.SaveAs(os.path.abspath(output_file), FileFormat=16)
    doc.Close()
    word.Quit()
    time.sleep()
    return output_file


class ImgToImg:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.ext = [mf.Return(mf.getExtension(ww)).getValues(0).lower() for ww in self.file]
        self.stop = False
        self.progressShow = None
    
    def toIMG(self, path, param={None:None}):
        param = param if type(param) == dict else {None:None}
        _ = list()
        for xx in range(len(self.file)):
            if not self.stop:
                if self.progressShow:
                    self.progressShow(int(xx * 100 / len(path)))
                try:
                    param = param if type(param) == dict else {None:param}
                    path = path if type(path) == list else [path]
                    clss = [mf.Return(mf.getExtension(path[ww])) for ww in range(len(path))]
                    ext = [clss[ww].getValues(0).lower() if clss[ww].getValues(0) in "cur,bmp,gif,jpeg,jpg,ico,png,svg,tiff,tif,webp".rsplit(",") else "png" for ww in range(len(clss))]
                    clss = [mf.Return(mf.getName(path[ww])) for ww in range(len(path))]
                    name = [clss[ww].getValues(0) for ww in range(len(clss))]
                    output = f"{TEMPCACHES}/tc-{dt.now()}.png".replace(":", "") if xx >= len(ext) else ".".join([name[xx], ext[xx]])
                    exec(f"""a=self.to{'PNG' if xx>=len(ext) else ext[xx].upper()}('{self.file[xx].replace("\\", "/")}', '{output.replace("\\", "/")}', param)
_.append(mf.setReturn(True, [type(a)], [a]))""")
                    if not _[-1]: del(_[-1])
                except Exception as e:
                    _.append(mf.setReturn(False, [type(e)], [e]))
                if self.progressShow:
                    self.progressShow(int(xx * 100 / len(path)))
            else:
                for i in _:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                _=[]
                break
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _
    
    def toPNG(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        if mf.Return(mf.getExtension(pathF)).getValues(0) == 'svg':
            with WandImage(filename=pathF) as img:
                img.format = 'png'
                img.save(filename=mf.Return(mf.getName(pathT)).getValues(0) + ".png", format="PNG")
        else:
            with Image.open(pathF) as img:
                output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".png"
                try: img.convert("RGBA")
                except: pass
                img.save(output_file, format="PNG")
        return pathT
    
    def toCUR(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        a = self.toICO(pathF, f"{TEMPCACHES}\\any-png-curany-{dt.now()}.png".replace(':', '.'))
        output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".cur"
        win32api.CopyFile(a, output_file, True)
        os.remove(a)
        return output_file
    
    def toBMP(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        a = self.toPNG(pathF, f"{TEMPCACHES}\\any-png-bmpany-{dt.now()}.png".replace(':', '.'))
        output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".bmp"
        with Image.open(a) as img:
            img.save(output_file, format="BMP")
        os.remove(a)
        return output_file
    
    def toGIF(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        a = self.toPNG(pathF, f"{TEMPCACHES}\\any-png-gifany-{dt.now()}.png".replace(':', '.'))
        output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".gif"
        with Image.open(a) as img:
            img.save(output_file, format="GIF")
        os.remove(a)
        return output_file
    
    def toICO(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        a = self.toPNG(pathF, f"{TEMPCACHES}\\any-png-icoany-{dt.now()}.png".replace(':', '.'))
        output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".ico"
        with Image.open(a) as img:
            img.save(output_file, format="ICO")
        os.remove(a)
        return output_file
    
    def toJPEG(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        a = self.toPNG(pathF, f"{TEMPCACHES}\\any-png-jpegany-{dt.now()}.png".replace(':', '.'))
        output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".jpeg"
        with Image.open(a) as img:
            img.save(output_file, format="JPEG")
        os.remove(a)
        return output_file
    
    def toJPG(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        a = self.toPNG(pathF, f"{TEMPCACHES}\\any-png-jpgany-{dt.now()}.png".replace(':', '.'))
        output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".jpg"
        with Image.open(a) as img:
            img.save(output_file, format="JPEG")
        os.remove(a)
        return output_file
    
    def toSVG(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        a = self.toPNG(pathF, f"{TEMPCACHES}\\any-png-svgany-{dt.now()}.png".replace(':', '.'))
        output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".svg"
        with WandImage(filename=a, resolution=param.get("resolution", 1080)) as img:
            img.format = 'svg'
            img.save(filename=output_file)
        os.remove(a)
        return output_file
    
    def toTIFF(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        a = self.toPNG(pathF, f"{TEMPCACHES}\\any-png-tiffany-{dt.now()}.png".replace(':', '.'))
        output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".tiff"
        with Image.open(a) as img:
            img.save(output_file, format="TIFF")
        os.remove(a)
        return output_file

    def toTIF(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        a = self.toPNG(pathF, f"{TEMPCACHES}\\any-png-tifany-{dt.now()}.png".replace(':', '.'))
        output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".tif"
        with Image.open(a) as img:
            img.save(output_file, format="TIFF")
        os.remove(a)
        return output_file
    
    def toWEBP(self, pathF, pathT, param={None:None}):
        param = param if type(param) == dict else {None:None}
        a = self.toPNG(pathF, f"{TEMPCACHES}\\any-png-webpany-{dt.now()}.png".replace(':', '.'))
        output_file = mf.Return(mf.getName(pathT)).getValues(0) + ".webp"
        with Image.open(a) as img:
            img.save(output_file, format="WEBP")
        os.remove(a)
        return output_file

    def kill(self):
        self.stop = True
    
    
class VideoToImg:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.ext = [mf.Return(mf.getExtension(ww)).getValues(0).lower() for ww in self.file]
        self.stop = False
        self.toKill = None
        self.progressShow = None
    
    def Convert(self, pathF, pathT, interval=None):
        _ = list()
        for ww in range(len(pathF)):
            _.append(os.path.join(TEMPCACHES, f"tc-videomg-{ww}-{dt.now()}".replace(":", "")))
            os.mkdir(_[ww])
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-videomg-{dt.now()}".replace(":", "")))
            os.mkdir(_[-1])
        __ = list()
        x, y = 0, 0
        if self.progressShow:
            self.progressShow(x)
        for ww in range(len(pathF)):
            if not self.stop:
                y = 0
                try:
                    vidcap = cv2.VideoCapture(pathF[ww])
                    if not vidcap.isOpened():
                        __.append(mf.setReturn(False, [type(pathF[ww])], [pathF[ww]]))
                        return
                    fps = vidcap.get(cv2.CAP_PROP_FPS)
                    inter = [interval[ww][0] * fps if interval[ww][0] * fps <= int(vidcap.get(cv2.CAP_PROP_FRAME_COUNT)) else 0, interval[ww][1] * fps if interval[ww][1] * fps <= int(vidcap.get(cv2.CAP_PROP_FRAME_COUNT)) else int(vidcap.get(cv2.CAP_PROP_FRAME_COUNT))] if interval else [0, int(vidcap.get(cv2.CAP_PROP_FRAME_COUNT))]
                    for xx in range(int(inter[1] - inter[0])):
                        if not self.stop:
                            vidcap.set(cv2.CAP_PROP_POS_FRAMES, inter[0] + xx)
                            success, image = vidcap.read()
                            if success:
                                output = f"{_[ww]}/tc-{dt.now()}.png".replace(":", "")
                                cv2.imwrite(output, image)
                                y = int(xx * 100 / (inter[1] - inter[0]))
                                if self.progressShow:
                                    self.progressShow(y)
                        else:break
                    __.append(mf.setReturn(True, [type(_[ww])], [_[ww]]))
                    x = int(ww * 100 / len(pathF))
                    vidcap.release()
                except Exception as e:
                    __.append(mf.setReturn(False, [type(e)], [e]))
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __
    
    def toIMG(self, path, extension="png", interval=None, modify={}):
        temp = self.Convert(self.file, [], interval)
        extension = extension if type(extension) == list else [extension]
        path = path if type(path) == list else [path]
        _ = list()
        for pp, p in enumerate(path):
            if not self.stop:
                if not os.path.exists(p):os.mkdir(p)
                p.replace("\\", "/")
                self.toKill = ModifyToImg([f"{mf.Return(temp[pp]).getValues(0)}/{ww}".replace("\\", "/") for ww in os.listdir(mf.Return(temp[pp]).getValues(0))])
                _.append(self.toKill.toIMG([f"{p}/{ww+1}.{extension[pp]}" for ww in range(len(os.listdir(mf.Return(temp[pp]).getValues(0))))], modify))
            else:
                for i in _:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                _=[]
                break
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop=True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)


class TextToImg:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None
    
    def Convert(self, pathF, pathT, bgColor="white", fgColor="black", width=800, height=600, margin=10, font=["arial", 20]):
        pathT
        _ = list()
        for ww in range(len(pathF)):
            _.append(os.path.join(TEMPCACHES, f"tc-textmg-{ww}-{dt.now()}.png".replace(":", "")))
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-textmg-{dt.now()}.png".replace(":", "")))
        __ = list()
        fx = 0
        for ww in range(len(pathF)):
            if not self.stop:
                try:
                    fy = 0
                    if self.progressShow:
                        self.progressShow(int(fy))
                    with open(pathF[ww], 'r', encoding="utf-8") as file:
                        text = file.read()
                    img = Image.new('RGB', (width, height), color=bgColor)
                    draw = ImageDraw.Draw(img)
                    try:
                        font_family = font[0]
                        font_size = font[1]
                        font = ImageFont.truetype(f"{font_family}.ttf", font_size)
                    except:
                        font = ImageFont.load_default()
                    x, y = margin, margin
                    for line in text.split('\n'):
                        draw.text((x, y), line, font=font, fill=fgColor)
                        y += font.getbbox(line)[3] + 2
                        fy = (text.split("\n").index(line) + 1) * 100 / len(text.split("\n"))
                        if self.progressShow:
                            self.progressShow(int(fy))
                    img.save(_[ww], 'PNG')
                    fx = int(ww * 100 / len(pathF))
                    fy = 100
                    __.append(mf.setReturn(True, [type(_[ww])], [_[ww]]))
                except Exception as e:
                    __.append(mf.setReturn(False, [type(e)], [e]))
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __
    
    def toIMG(self, path, size=[997, 1080, 10], color=["white", "black"], font=["arial", 20], modify={}):
        temp = self.Convert(self.file, [], color[0], color[1], size[0], size[1], size[2], font)
        path = path if type(path) == list else [path]
        _ = list()
        for ww in range(len(temp)):
            if not self.stop:
                if mf.Return(temp[ww]).getSucced():
                    inp = fr'{mf.Return(temp[ww]).getValues(0)}'.replace("\\", "/")
                    outp = fr'{path[ww]}'.replace("\\", "/")
                    self.toKill = ModifyToImg(inp)
                    _.append(self.toKill.toIMG(outp, modify)[0])
                    try:
                        os.remove(mf.Return(temp[ww]).getValues(0))
                    except Exception as e:
                        print(e)
                else:
                    print(mf.Return(temp[ww]).getValues())
            else:
                for i in _:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                _=[]
                break
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)


class PdfToImg:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None
    
    def Convert(self, pathF, pathT):
        _ = list()
        for ww in range(len(pathF)):
            _.append(os.path.join(TEMPCACHES, f"tc-pdfmg-{ww}-{dt.now()}".replace(":", "")))
            os.mkdir(_[ww])
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-pdfmg-{dt.now()}".replace(":", "")))
            os.mkdir(_[-1])
        __ = list()
        x, y = 0, 0
        if self.progressShow:
            self.progressShow(x)
        for ww in range(len(pathF)):
            if not self.stop:
                y = 0
                try:
                    images = fitz.open(pathF[ww])
                    for xx, image in enumerate(images):
                        page = images.load_page(xx)
                        pix = page.get_pixmap()
                        image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        image_path = f"{_[ww]}\\Page{xx+1}.png".replace("\\", "/")
                        image.save(image_path, 'PNG')
                        y = int(xx * 100 / len(images))
                        if self.progressShow:
                            self.progressShow(y)
                    __.append(mf.setReturn(True, [type(_[ww])], [_[ww]]))
                    x = int(ww * 100 / len(pathF))
                    y = 100
                except Exception as e:
                    __.append(mf.setReturn(False, [type(e)], [e]))
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        print(__)
        return __
    
    def toIMG(self, path, extension="png", modify={}):
        temp = self.Convert(self.file, [])
        extension = extension if type(extension) == list else [extension]
        path = path if type(path) == list else [path]
        _ = list()
        for pp, p in enumerate(path):
            if not self.stop:
                if not os.path.exists(p):os.mkdir(p)
                p.replace("\\", "/")
                self.toKill = ModifyToImg([f"{mf.Return(temp[pp]).getValues(0)}/{ww}".replace("\\", "/") for ww in os.listdir(mf.Return(temp[pp]).getValues(0))])
                _.append(self.toKill.toIMG([f"{p}/{ww+1}.{extension[pp]}" for ww in range(len(os.listdir(mf.Return(temp[pp]).getValues(0))))], modify))
            else:
                for i in _:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                _=[]
                break
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)
    
           
class WordToImg:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None
    
    def Convert(self, pathF, pathT):
        pathT
        _ = list()
        for ww in range(len(pathF)):
            type(ww)
            _.append(os.path.join(TEMPCACHES, f"tc-wordmg-{ww}-{dt.now()}.pdf".replace(":", "")))
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-wordmg-{dt.now()}.pdf".replace(":", "")))
        __ = list()
        for xx, ww in enumerate(pathF):
            if not self.stop:
                if self.progressShow:
                    self.progressShow(int(xx * 100 / len(_)))
                try:
                    temp = wtw(pathF[xx], f"{TEMPCACHES}/tc-{dt.now()}.docx".replace(":", ""))
                    cwp(temp, _[xx])
                    __.append(mf.setReturn(True, [type(_[xx])], [_[xx]]))
                except Exception as e:
                    __.append(mf.setReturn(False, [type(e)], [e]))
                if self.progressShow:
                    self.progressShow(int(xx * 100 / len(_)))
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __
    
    def toIMG(self, path, extension="png", modify={}):
        temp = self.Convert(self.file, [])
        print(temp)
        path = path if type(path) == list else [path]
        extension = extension if type(extension) == list else [extension]
        self.toKill = PdfToImg([mf.Return(ww).getValues(0) for ww in temp])
        _ = self.toKill.toIMG(path, extension, modify)
        for ww in temp:
            if not self.stop:
                try:
                    os.remove(mf.Return(ww).getValues(0))
                except Exception as e:
                    print(e)
            else:
                for i in _:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                _=[]
                break
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)


class PpToImg:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None

    def Convert(self, pathF, pathT):
        pathT
        _ = list()
        for ww in range(len(pathF)):
            _.append(os.path.join(TEMPCACHES, f"tc-ppmgc-{ww}-{dt.now()}".replace(":", "")))
            os.mkdir(_[ww])
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-ppmgc-{dt.now()}".replace(":", "")))
            os.mkdir(_[-1])
        __ = list()
        x, y = 0, 0
        if self.progressShow:
            self.progressShow(x)
        for ww in range(len(pathF)):
            if not self.stop:
                y = 0
                try:
                    prs = Presentation(pathF[ww])
                    for slide_number, slide in enumerate(prs.slides):
                        img = Image.new('RGB', (960, 720), color='white')
                        img.save(f"{_[ww]}/slide{slide_number + 1}.png")
                        type(slide)
                        y = int(slide_number * 100 / len(prs.slides))
                        if self.progressShow:
                            self.progressShow(y)
                    __.append(mf.setReturn(True, [type(_[ww])], [_[ww]]))
                    x = int(ww * 100 / len(pathF))
                    y = 100
                except Exception as e:
                    __.append(mf.setReturn(False, [type(e)], [e]))
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __

    def toIMG(self, path, extension="png", modify={}):
        temp = self.Convert(self.file, [])
        extension = extension if type(extension) == list else [extension]
        path = path if type(path) == list else [path]
        _ = list()
        for pp, p in enumerate(path):
            if not self.stop:
                if not os.path.exists(p):os.mkdir(p)
                p.replace("\\", "/")
                self.toKill = ModifyToImg([f"{mf.Return(temp[pp]).getValues(0)}/{ww}".replace("\\", "/") for ww in os.listdir(mf.Return(temp[pp]).getValues(0))])
                _.append(self.toKill.toIMG([f"{p}/{ww+1}.{extension[pp]}" for ww in range(len(os.listdir(mf.Return(temp[pp]).getValues(0))))], modify))
            else:
                for i in _:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                _=[]
                break
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)


class LinkToImg:

    def __init__(self, link):
        file = link
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None
    
    def Convert(self, pathF, pathT, qr={}):
        pathT
        qr = {"pixel":qr.get("pixel", 10),
            "margin":qr.get("margin", 30),
            "bg":qr.get("bg", "white"),
            "fg":qr.get("fg", "black"),
            "version":qr.get("version", 1),
            "correction":qr.get("correction", "L"),
            "width":qr.get("width", 300),
            "height":qr.get("height", 300),
            "image":qr.get("image", None),
            "impact":qr.get("impact", .2),
            "reduction":qr.get("reduction", .5)}
        _ = list()
        for ww in range(len(pathF)):
            _.append(os.path.join(TEMPCACHES, f"tc-linkc-{ww}-{dt.now()}.png".replace(":", "")))
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-linkc-{dt.now()}.png".replace(":", "")))
        __ = list()
        x, y = 0, 0
        if self.progressShow:
            self.progressShow(x)
        for xx, ww in enumerate(pathF):
            if not self.stop:
                y = 0
                if self.progressShow:
                    self.progressShow(x)
                try:
                    correction_levels = {'L': qrcode.constants.ERROR_CORRECT_L,
                                        'M': qrcode.constants.ERROR_CORRECT_M,
                                        'Q': qrcode.constants.ERROR_CORRECT_Q,
                                        'H': qrcode.constants.ERROR_CORRECT_H}
                    qrc = qrcode.QRCode(
                        version=qr["version"],
                        error_correction=correction_levels.get(qr["correction"], qrcode.constants.ERROR_CORRECT_L),
                        box_size=qr["pixel"],
                        border=qr["margin"] // qr["pixel"],
                    )
                    qrc.add_data(ww)
                    qrc.make(fit=True)
                    img = qrc.make_image(fill=qr["fg"], back_color=qr["bg"])
                    img = img.resize((qr["width"], qr["height"]))
                    if qr["image"]:
                        logo = Image.open(qr["image"])
                        if logo.mode != 'RGBA':
                            logo = logo.convert('RGBA')
                        logo_size = int(min(img.size) * qr["impact"])
                        logo = logo.resize((logo_size, logo_size))
                        logo_pos = ((img.width - logo.width) // 2, (img.height - logo.height) // 2)
                        img.paste(logo, logo_pos, logo)
                    if qr["reduction"] != 1:
                        img = img.resize((int(img.width * qr["reduction"]), int(img.height * qr["reduction"])))
                    img.save(_[xx])
                    __.append(mf.setReturn(True, [type(_[xx])], [_[xx]]))
                except Exception as e:
                    __.append(mf.setReturn(True, [type(e)], [e]))
                x = int(xx * 100 / len(pathF))
                y = 100
                if self.progressShow:
                    self.progressShow(x)
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __

    def toIMG(self, path, qr={}):
        temp = self.Convert(self.file, [], qr)
        path = path if type(path) == list else [path]
        _ = list()
        for ww in range(len(temp)):
            if not self.stop:
                if mf.Return(temp[ww]).getSucced():
                    inp = fr'{mf.Return(temp[ww]).getValues(0)}'.replace("\\", "/")
                    outp = fr'{path[ww]}'.replace("\\", "/")
                    self.toKill = ImgToImg(inp)
                    _.append(self.toKill.toIMG(outp, {None:None})[0])
                    try:
                        os.remove(mf.Return(temp[ww]).getValues(0))
                    except Exception as e:
                        print(e)
                else:
                    print(mf.Return(temp[ww]).getValues())
            else:
                for i in _:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                _=[]
                break
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)


class ModifyToImg:

    def __init__(self, file):
        self.file = file if type(file) == list else [file]
        self.stop = False
        self.toKill = None
        self.progressShow = None
    
    def Convert(self, pathF, pathT, modify={}):
        pathT
        new = {"width":modify.get("width", None),
            "height":modify.get("height", None),
            "scale":modify.get("scale", 1.0),
            "invertcolor":modify.get("invertColor", False),
            "mirror": modify.get("miror", False),
            "century":modify.get("century", False),
            "rounded":modify.get("rounded", 0),
            "replace":modify.get("replace", None),
            "angle":modify.get("angle", 0),
            "rognage":modify.get("rognage", (0, 0, 100, 100)),
            "brightness":modify.get("brightness", 0),
            "saturation":modify.get("saturation", 0),
            "stickers":modify.get("stickers", None),
            "deleteColor":modify.get("deleteColor", None)}
        """
            Convert an image with the specified configurations.
            
            :param pathF: Source image path (PNG format).
            :param pathT: Destination image path (PNG format).
            :param width: Target width of the image (default: original width).
            :param height: Target height of the image (default: original height).
            :param scale: Scale factor for the image (default: 1.0).
            :param invertcolor: If True, invert the colors.
            :param miror: Apply a miror effect to imgae
            :param century: If True, convert the image to black and white.
            :param rounded: Border radius for rounded corners (0-100, default: 0).
            :param replace: List of color replacements [[old_color, new_color, smooth], ...] (default: None).
            :param angle: Rotation angle for the image (default: 0°).
            :param rognage: Crop region (left, upper, right, lower) in percentages (default: (0, 0, 100, 100)).
            :param brightness: Brightness adjustment (-100 to 100, default: 0).
            :param saturation: Saturation adjustment (-100 to 100, default: 0).
            :param stickers: List of stickers [[path, [x, y], [width, height], angle], ...] (default: None).
            :param deleteColor: List of colors to make transparent [[color, smooth], ...] (default: None).
            """
        _ = list()
        for ww in range(len(pathF)):
            _.append(os.path.join(TEMPCACHES, f"tc-{ww}-{dt.now()}.png".replace(":", "")))
        while(len(pathF) > len(_)):
            _.append(os.path.join(TEMPCACHES, f"tc-mdc-{dt.now()}.png".replace(":", "")))
        __ = list()
        fx, fy = 0, 0
        if self.progressShow:
            self.progressShow(fx)
        self.toKill = ImgToImg(pathF)
        pathF = self.toKill.toIMG([f"{TEMPCACHES}/tc-imgc-{ww}-{dt.now()}.png".replace(":", "") for ww in range(len(pathF))], {None:None})
        for xx, ww in enumerate(pathF):
            if not self.stop:
                fy = 0
                if mf.Return(ww).getSucced():
                    try:
                        # Load the source image
                        img = Image.open(mf.Return(ww).getValues(0))
                        try:img.convert("RGBA")
                        except Exception as e: print(e)
                        orig_width, orig_height = img.size
                        pixels = img.load()

                        # Apply scaling
                        if new["scale"] != 1.0:
                            img = img.resize((int(orig_width * new["scale"]), int(orig_height * new["scale"])))
                        
                        # Apply width and height adjustments
                        target_width = new["width"] if new["width"] else img.width
                        target_height = new["height"] if new["height"] else img.height
                        img = img.resize((target_width, target_height))

                        # applying miror effect
                        if new["mirror"]:
                            img = img.transpose(Image.FLIP_LEFT_RIGHT)
                        
                        # Apply cropping
                        left = new["rognage"][0] / 100 * img.width
                        upper = new["rognage"][1] / 100 * img.height
                        right = new["rognage"][2] / 100 * img.width
                        lower = new["rognage"][3] / 100 * img.height
                        img = img.crop((left, upper, right, lower))
                        
                        # Apply color inversion
                        if new["invertcolor"]:
                            img = Image.fromarray(255 - np.array(img))

                        # Convert to black and white if `century` is True
                        if new["century"]:
                            img = img.convert("L").convert("RGBA")

                        # Adjust brightness
                        if new["brightness"] != 0:
                            enhancer = ImageEnhance.Brightness(img)
                            img = enhancer.enhance(1 + new["brightness"] / 100)

                        # Adjust saturation
                        if new["saturation"] != 0:
                            enhancer = ImageEnhance.Color(img)
                            img = enhancer.enhance(1 + new["saturation"] / 100)

                        # Replace colors with smooth tolerance
                        if new["replace"]:
                            data = np.array(img)
                            for old_color, new_color, smooth in new["replace"]:
                                r1, g1, b1, a1 = old_color
                                r2, g2, b2, a2 = new_color
                                # diff = smooth / 255  # Normalize smooth
                                mask = (
                                    (np.abs(data[:,:, 0] - r1) <= smooth) & 
                                    (np.abs(data[:,:, 1] - g1) <= smooth) & 
                                    (np.abs(data[:,:, 2] - b1) <= smooth) & 
                                    (np.abs(data[:,:, 3] - a1) <= smooth)
                                )
                                data[mask] = [r2, g2, b2, a2]
                            img = Image.fromarray(data)

                        # Remove colors with smooth tolerance
                        if new["deleteColor"]:
                            if new["deleteColor"] == "all":
                                # Supprime toutes les couleurs : image transparente
                                img = Image.new("RGBA", img.size, (0, 0, 0, 0))
                            else:
                                for y in range(img.height):
                                    for x in range(img.width):
                                        r, g, b, a = pixels[x, y]
                                        for color, smooth in new["deleteColor"]:
                                            type(a)
                                            cr, cg, cb, _ = mf.getColor(color)
                                            # Calcul de la distance entre les couleurs
                                            if abs(r - cr) <= smooth and abs(g - cg) <= smooth and abs(b - cb) <= smooth:
                                                pixels[x, y] = (0, 0, 0, 0)  # Transparence complète

                        # Rotate the image
                        if new["angle"] != 0:
                            img = img.rotate(new["angle"], expand=True)

                        # Add rounded corners
                        if new["rounded"] > 0:
                            radius = int(min(img.size) * (new["rounded"] / 100))
                            mask = Image.new("L", img.size, 0)
                            draw = ImageDraw.Draw(mask)
                            draw.rounded_rectangle([(0, 0), img.size], radius, fill=255)
                            img.putalpha(mask)

                        # Add stickers
                        if new["stickers"]:
                            for sticker in new["stickers"]:
                                sticker_path, (x, y), (sticker_width, sticker_height), sticker_angle = sticker
                                sticker_img = Image.open(sticker_path).convert("RGBA")
                                sticker_img = sticker_img.resize((sticker_width, sticker_height))
                                sticker_img = sticker_img.rotate(sticker_angle, expand=True)
                                x_pos = int(x / 100 * img.width - sticker_width / 2)
                                y_pos = int(y / 100 * img.height - sticker_height / 2)
                                img.paste(sticker_img, (x_pos, y_pos), sticker_img)

                        # Save the final image
                        img.save(_[xx], format="PNG")
                        __.append(mf.setReturn(True, [type(_[xx])], [_[xx]]))
                    except Exception as e:
                        __.append(mf.setReturn(False, [type(e)], [e]))
                else:
                    __.append(ww)
                try:
                    os.remove(mf.Return(ww).getValues(0))
                except Exception as e:
                    print(e)
                y = 100
                x = int(xx * 100 / len(pathF))
                if self.progressShow:
                    self.progressShow(x)
            else:
                for i in __:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                break
        if self.progressShow:
            self.progressShow(100)
        return __

    def toIMG(self, path, modify={}):
        temp = self.Convert(self.file, [], modify)
        path = path if type(path) == list else [path]
        _ = list()
        for ww in range(len(temp)):
            if not self.stop:
                if mf.Return(temp[ww]).getSucced():
                    inp = fr'{mf.Return(temp[ww]).getValues(0)}'.replace("\\", "/")
                    outp = fr'{path[ww]}'.replace("\\", "/")
                    self.toKill = ImgToImg(inp)
                    _.append(self.toKill.toIMG(outp, {None:None})[0])
                else:
                    _.append(temp[ww])
            else:
                for i in _:
                    try:
                        os.remove(mf.Return(i).getValues(0))
                    except:pass
                _=[]
                break
        if self.progressShow:
            self.progressShow(100, True, self.stop, _)
        return _

    def kill(self):
        self.stop = True
        if self.toKill:
            try:
                self.toKill.kill()
            except Exception as e:
                print(e)


